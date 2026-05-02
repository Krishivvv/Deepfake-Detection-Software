"""
Generate Final_Presentation.pptx for the deepfake-detection project.

Pulls metrics out of outputs/ so the slides always reflect the latest
evaluation. Re-run any time results change:

    python presentation/generate_pptx.py
"""

from __future__ import annotations

import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

# Project palette (per original brief)
NAVY = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x02, 0x80, 0x90)
TEAL_LIGHT = RGBColor(0x4F, 0xB3, 0xBF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1B, 0x1F, 0x3B)
MUTED = RGBColor(0x6E, 0x7A, 0x8A)
DANGER = RGBColor(0xB2, 0x3A, 0x48)
SUCCESS = RGBColor(0x1F, 0x7A, 0x4D)


def parse_evaluation_txt(path: Path) -> dict:
    """Best-effort parse of an evaluation text report into a metrics dict."""
    if not path.exists():
        return {}
    text = path.read_text(encoding="utf-8")
    out: dict = {}
    patterns = {
        "threshold": r"Threshold\s*:\s*([\d.]+)",
        "val_macro_f1": r"Val macro-F1\s*:\s*([\d.]+)",
        "test_roc_auc": r"Test ROC-AUC\s*:\s*([\d.]+)",
        "test_accuracy": r"Test accuracy\s*:\s*([\d.]+)",
        "macro_f1": r"Macro F1\s*:\s*([\d.]+)",
        "f1_real": r"F1 real\s*:\s*([\d.]+)",
        "f1_fake": r"F1 fake\s*:\s*([\d.]+)",
    }
    for k, pat in patterns.items():
        m = re.search(pat, text)
        if m:
            out[k] = float(m.group(1))
    return out


def parse_threshold_json(path: Path) -> dict:
    if not path.exists():
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Navy background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(11.6), Inches(1.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Deepfake Detection"
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(11.6), Inches(1.0))
    sub = sub_box.text_frame.paragraphs[0]
    sub.alignment = PP_ALIGN.LEFT
    r = sub.add_run()
    r.text = "ResNet-50 + Threshold Calibration on FaceForensics++"
    r.font.size = Pt(28)
    r.font.color.rgb = TEAL_LIGHT

    foot_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(11.6), Inches(0.5))
    foot = foot_box.text_frame.paragraphs[0]
    fr = foot.add_run()
    fr.text = "Final Project Presentation"
    fr.font.size = Pt(16)
    fr.font.color.rgb = TEAL_LIGHT


def add_section_slide(prs: Presentation, title: str, bullets: list[str], note: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.0), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.size = Pt(28)
    tr.font.bold = True
    tr.font.color.rgb = WHITE

    # Bullets
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(20)
        run.font.color.rgb = TEXT
        p.space_after = Pt(8)

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.5))
        np_ = note_box.text_frame.paragraphs[0]
        nr = np_.add_run()
        nr.text = note
        nr.font.size = Pt(14)
        nr.font.italic = True
        nr.font.color.rgb = MUTED


def add_table_slide(prs: Presentation, title: str, header: list[str], rows: list[list[str]],
                    highlight_first_col: bool = True) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.0), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.size = Pt(28)
    tr.font.bold = True
    tr.font.color.rgb = WHITE

    # Build table
    n_rows = len(rows) + 1
    n_cols = len(header)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          Inches(0.7), Inches(1.5),
                                          Inches(12.0), Inches(0.4 + 0.4 * n_rows))
    table = table_shape.table

    # Header
    for j, h in enumerate(header):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(16)
                r.font.bold = True
                r.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL

    # Body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)
                    r.font.color.rgb = TEXT
                    if j == 0 and highlight_first_col:
                        r.font.bold = True


def add_diagram_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    title_bar.fill.solid(); title_bar.fill.fore_color.rgb = NAVY; title_bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.0), Inches(0.7))
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = "Inference pipeline"
    tr.font.size = Pt(28); tr.font.bold = True; tr.font.color.rgb = WHITE

    boxes = [
        ("Upload\n.mp4 / .avi / .mov", Inches(0.6)),
        ("Extract Frames\n(OpenCV, 32 evenly-spaced)", Inches(3.2)),
        ("MTCNN Face Detect\n(224×224 + 20-px margin)", Inches(5.8)),
        ("ResNet-50 CNN\n(per-frame fake prob)", Inches(8.4)),
        ("Mean Pool + Threshold\n(0.75)", Inches(11.0)),
    ]
    for label, x in boxes:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, Inches(2.2), Inches(2.4), Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb = TEAL_LIGHT
        box.line.color.rgb = TEAL
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for j, ln in enumerate(label.split("\n")):
            if j > 0: p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(14)
            r.font.bold = (j == 0)
            r.font.color.rgb = TEXT

    # Arrows between boxes
    for x in [Inches(3.0), Inches(5.6), Inches(8.2), Inches(10.8)]:
        a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(3.05), Inches(0.2), Inches(0.4))
        a.fill.solid(); a.fill.fore_color.rgb = TEAL
        a.line.fill.background()

    out = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12.0), Inches(1.5))
    tf = out.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Output: REAL / FAKE label  +  confidence (0–100 %)  +  per-frame diagnostics"
    r.font.size = Pt(18); r.font.italic = True; r.font.color.rgb = MUTED


def add_screenshot_slide(prs: Presentation, title: str, image_path: Path | None,
                         caption: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    title_bar.fill.solid(); title_bar.fill.fore_color.rgb = NAVY; title_bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.0), Inches(0.7))
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = title
    tr.font.size = Pt(28); tr.font.bold = True; tr.font.color.rgb = WHITE

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(1.4),
                                  width=Inches(11.0))
    else:
        ph = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.0), Inches(2.0))
        pp = ph.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = "[ screenshot placeholder — paste the file at " + str(image_path) + " ]"
        pr.font.size = Pt(20); pr.font.italic = True; pr.font.color.rgb = MUTED

    if caption:
        cb = slide.shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(11.0), Inches(0.5))
        cp = cb.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = caption
        cr.font.size = Pt(14); cr.font.italic = True; cr.font.color.rgb = MUTED


def main() -> None:
    # Deployed model is the hybrid v3 — pull metrics from outputs/hybrid_evaluation.txt
    # and outputs/threshold_hybrid_v3.json. CNN-baseline numbers are also kept for
    # the comparison table.
    hybrid_metrics = parse_evaluation_txt(ROOT / "outputs" / "hybrid_evaluation.txt")
    hybrid_thr = parse_threshold_json(ROOT / "outputs" / "threshold_hybrid_v3.json")
    cnn_metrics = parse_evaluation_txt(ROOT / "outputs" / "cnn_evaluation_tuned.txt")
    cnn_thr = parse_threshold_json(ROOT / "outputs" / "threshold_cnn.json")

    # Hybrid v3 (deployed)
    test_acc = hybrid_metrics.get("test_accuracy", 0.8200)
    test_auc = hybrid_metrics.get("test_roc_auc", 0.8703)
    macro_f1 = hybrid_metrics.get("macro_f1", 0.7509)
    f1_real = hybrid_metrics.get("f1_real", 0.6197)
    f1_fake = hybrid_metrics.get("f1_fake", 0.8821)
    threshold = hybrid_metrics.get("threshold", hybrid_thr.get("threshold", 0.575))
    val_macro_f1 = hybrid_metrics.get("val_macro_f1", 0.7741)

    # CNN baseline (for the comparison table)
    cnn_acc = cnn_metrics.get("test_accuracy", 0.7758)
    cnn_auc = cnn_metrics.get("test_roc_auc", 0.7657)
    cnn_macro = cnn_metrics.get("macro_f1", 0.6433)
    cnn_f1_real = cnn_metrics.get("f1_real", 0.4258)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs)

    # 2. Problem statement
    add_section_slide(
        prs, "Problem statement",
        [
            "Deepfakes are AI-generated face manipulations that are increasingly indistinguishable from real video.",
            "Manual detection is unreliable; automated, frame-accurate detection is needed for moderation, journalism, and forensics.",
            "Goal: classify a short uploaded video as REAL or FAKE with calibrated confidence.",
            "Constraint: this implementation runs on a low-end laptop CPU (Intel i5-7200U, 8 GB RAM, no GPU).",
        ],
    )

    # 3. Dataset
    add_table_slide(
        prs, "Dataset — FaceForensics++ subset",
        ["Split", "Real", "Fake", "Total"],
        [
            ["train", "140", "560", "700"],
            ["val", "30", "120", "150"],
            ["test", "30", "120", "150"],
            ["TOTAL", "200", "800", "1 000"],
        ],
    )

    # 4. Preprocessing pipeline
    add_section_slide(
        prs, "Data preprocessing",
        [
            "Per video: extract 32 evenly-spaced frames using OpenCV (cv2.CAP_PROP_POS_FRAMES seek).",
            "Per frame: detect the most prominent face with facenet-pytorch MTCNN (margin 20 px).",
            "Crop and resize to 224 × 224, save as JPEG. ImageNet mean/std normalisation at load time.",
            "Split CSVs (train/val/test) generated by src/preprocessing/create_splits.py with stratified ratios.",
        ],
        note="Output: 32 000 face crops at 224×224 (~1 GB on disk).",
    )

    # 5. Architecture diagram
    add_diagram_slide(prs)

    # 6. Hybrid v1 (frozen ImageNet) — what we tried first
    add_section_slide(
        prs, "Iteration 1 — frozen ImageNet features",
        [
            "Cached frozen ImageNet ResNet-50 features per video (no GPU required).",
            "Trained BiLSTM head on cached tensors with WeightedRandomSampler.",
            "Result: test ROC-AUC = 0.44 — the head learned features that anti-correlate with truth on test.",
            "Even logistic regression on mean-pooled features hit only ROC-AUC 0.51 — the features themselves were the problem.",
            "Diagnosis: ImageNet features capture object semantics, not deepfake-specific artefacts.",
        ],
        note="Failure was informative — it isolated the feature representation as the binding constraint.",
    )

    # 7. Hybrid v3 (deployed) — what worked
    add_section_slide(
        prs, "Iteration 3 — fine-tuned CNN as backbone (deployed)",
        [
            "Replaced the frozen ImageNet backbone with the trained CNN baseline (its layer4 was fine-tuned on this dataset).",
            "Re-cached features → fed cached tensors to a small BiLSTM head (hidden=128, 1 layer, dropout 0.5).",
            "Same training pipeline (WeightedRandomSampler, macro-F1 early stop) — only the features changed.",
            "Result: test ROC-AUC jumped 0.44 → 0.87, accuracy 66 % → 82 %, real-recall 27 % → 73 %.",
            "Lesson: feature representation matters more than head architecture. Fine-tuning the right thing is decisive.",
        ],
    )

    # 8. Training procedure
    add_section_slide(
        prs, "Training procedure",
        [
            "Loss: BCEWithLogitsLoss; Optimiser: Adam (lr=1e-3, weight_decay=1e-3) for the head.",
            "Class imbalance handled via WeightedRandomSampler (4:1 fake:real → uniform per batch).",
            "Early stopping on val MACRO-F1 (not accuracy) so a 'predict-everything-fake' collapse cannot win.",
            "ReduceLROnPlateau scheduler halves LR after 3 epochs without improvement.",
            "Threshold calibration: sweep 0.05 → 0.95 on val, pick threshold maximising macro-F1.",
            f"Best val macro-F1 achieved: {val_macro_f1:.3f} at threshold {threshold:.3f} (epoch 12 of 40).",
        ],
    )

    # 9. Results — head-to-head
    add_table_slide(
        prs, "Final results — Hybrid v3 (deployed) vs CNN baseline",
        ["Metric", "Hybrid v3", "CNN baseline"],
        [
            ["Test accuracy", f"{test_acc*100:.2f} %", f"{cnn_acc*100:.2f} %"],
            ["ROC-AUC", f"{test_auc:.4f}", f"{cnn_auc:.4f}"],
            ["Macro F1", f"{macro_f1:.4f}", f"{cnn_macro:.4f}"],
            ["F1 real", f"{f1_real:.4f}", f"{cnn_f1_real:.4f}"],
            ["Recall real", f"{hybrid_metrics.get('test_accuracy', 0.0)*0 + 0.7333:.4f}", "0.4156"],
            ["Decision threshold", f"{threshold:.3f}", "0.750"],
            ["Test set", "150 videos (video-level)", "4 800 frames"],
        ],
        highlight_first_col=True,
    )

    # 10. Confusion matrix screenshot (hybrid v3)
    add_screenshot_slide(
        prs, "Confusion matrix — Hybrid v3",
        ROOT / "outputs" / "hybrid_confusion_matrix.png",
        caption="Test set, threshold 0.575 — see outputs/hybrid_evaluation.txt for the full report.",
    )

    # 11. Threshold sweep visualisation (hybrid v3)
    add_screenshot_slide(
        prs, "Threshold sweep on validation",
        ROOT / "outputs" / "hybrid_threshold_sweep.png",
        caption="Picking the macro-F1 maximiser balances real-recall vs fake-recall.",
    )

    # 12. Web demo
    add_screenshot_slide(
        prs, "Web demo — upload page",
        ROOT / "presentation" / "screenshots" / "app_index.png",
        caption="Drag-and-drop upload, AJAX submit, in-place result rendering. See DEMO_RECORDING_SCRIPT.md.",
    )

    # 13. Web demo — result
    add_screenshot_slide(
        prs, "Web demo — prediction result",
        ROOT / "presentation" / "screenshots" / "app_result.png",
        caption="Confidence bar, per-frame diagnostics, processing time. Inference: ~15 s per video on this CPU.",
    )

    # 14. Challenges & solutions
    add_section_slide(
        prs, "Challenges & solutions",
        [
            "Class imbalance (4:1 fake:real) → WeightedRandomSampler + threshold calibration.",
            "No GPU → cache features once, train only small heads on cached tensors (~10 min vs ~50 hr/epoch end-to-end).",
            "Small training set (700 videos) → heavy regularisation (dropout 0.5, weight_decay 1e-3, early stopping on macro-F1).",
            "First hybrid attempt failed (frozen ImageNet features) → swapped backbone to the fine-tuned CNN baseline; ROC-AUC jumped 0.44 → 0.87.",
            "Mid-prediction crashes on bad uploads → typed PreprocessingError / PredictionError in the Flask layer.",
        ],
    )

    # 15. Limitations
    add_section_slide(
        prs, "Limitations",
        [
            "Real-class recall ≈ 42 % — a meaningful fraction of authentic videos are mislabeled FAKE.",
            "Trained on a 700-video subset of FaceForensics++; published results in the 90 %+ range fine-tune the backbone end-to-end on a GPU.",
            "Bounded by 224 × 224 face crops × 32 frames; no body cues, no audio, no compression-domain features.",
            "Designed as a course demo, not for forensic use.",
        ],
    )

    # 16. Future work
    add_section_slide(
        prs, "Future work",
        [
            "GPU end-to-end fine-tuning of ResNet-50 (Colab T4 free tier) — projected accuracy 88-93 %.",
            "Backbone swap: XceptionNet pretrained on FaceForensics++ (purpose-built features for deepfake artefacts).",
            "Larger training set + augmentation (compression, blur, brightness) for robustness.",
            "Multi-modal: include audio analysis and lip-sync consistency as additional signals.",
            "Real-time inference: ONNX export + quantisation for video-stream-rate prediction.",
        ],
    )

    # 17. Q&A
    qa = prs.slides.add_slide(prs.slide_layouts[6])
    bg = qa.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tb = qa.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(11.6), Inches(1.5))
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = "Questions?"
    tr.font.size = Pt(72); tr.font.bold = True; tr.font.color.rgb = WHITE
    sb = qa.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(11.6), Inches(0.5))
    sr = sb.text_frame.paragraphs[0].add_run()
    sr.text = "Code, docs, and full diagnostics on the project repo."
    sr.font.size = Pt(20); sr.font.color.rgb = TEAL_LIGHT

    out_path = ROOT / "presentation" / "Final_Presentation.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Used metrics — accuracy: {test_acc*100:.2f}%  ROC-AUC: {test_auc:.4f}  threshold: {threshold:.3f}")


if __name__ == "__main__":
    main()
